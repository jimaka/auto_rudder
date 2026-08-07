# -*- coding: utf-8 -*-
"""生成《船舶控制与节能算法调研》汇报 PPT（扩充版，20 页）"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

BASE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(BASE, "figures")

DARK = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x1F, 0x77, 0xB4)
RED = RGBColor(0xD6, 0x27, 0x28)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
ORANGE = RGBColor(0xE8, 0x7D, 0x0D)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Noto Sans CJK SC"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    return prs.slides.add_slide(BLANK)


def textbox(slide, x, y, w, h, lines, size=16, color=DARK, bold=False,
            align=PP_ALIGN.LEFT, line_spacing=1.15, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, dict overrides) 或纯字符串"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, str):
            line = (line, {})
        text, ov = line
        p.alignment = ov.get("align", align)
        p.line_spacing = ov.get("line_spacing", line_spacing)
        p.space_after = Pt(ov.get("space_after", 6))
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = ov.get("font", FONT)
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.color.rgb = ov.get("color", color)
    return tb


def header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.12))  # rectangle
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    textbox(slide, Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.9),
            [(title, {"size": 30, "bold": True, "color": DARK})])
    if subtitle:
        textbox(slide, Inches(0.52), Inches(1.0), Inches(12.3), Inches(0.5),
                [(subtitle, {"size": 14, "color": GRAY})])


def pic_fit(slide, path, x, y, max_w, max_h):
    """按比例放入图片，不超出 (max_w, max_h)，返回 shape"""
    from PIL import Image
    iw, ih = Image.open(path).size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    return slide.shapes.add_picture(path, x + int((max_w - w) / 2), y + int((max_h - h) / 2), w, h)


def bullets(slide, x, y, w, h, items, size=15):
    """items: list of (level, text, color|None, bold)"""
    lines = []
    for it in items:
        lvl, text = it[0], it[1]
        color = it[2] if len(it) > 2 and it[2] is not None else DARK
        bold = it[3] if len(it) > 3 else (lvl == 0)
        prefix = "▪ " if lvl == 0 else "    – "
        lines.append((prefix + text, {"size": size if lvl == 0 else size - 1.5,
                                      "color": color, "bold": bold,
                                      "space_after": 8 if lvl == 0 else 4}))
    return textbox(slide, x, y, w, h, lines)


def table(slide, x, y, w, h, data, col_widths=None, size=12, header_fill=BLUE):
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            gt.columns[i].width = Emu(int(w * cw / total))
    for r in range(rows):
        for c in range(cols):
            cell = gt.cell(r, c)
            cell.text = str(data[r][c])
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.line_spacing = 1.05
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(size if r else size + 0.5)
                    run.font.bold = (r == 0)
                    run.font.color.rgb = WHITE if r == 0 else DARK
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            elif r % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xEE, 0xF3, 0xFA)
    return gt


# ============ S01 封面 ============
s = add_slide()
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
textbox(s, Inches(1), Inches(2.2), Inches(11.3), Inches(1.6),
        [("船舶控制与节能算法调研汇报", {"size": 44, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER})])
textbox(s, Inches(1), Inches(3.7), Inches(11.3), Inches(1.2),
        [("MPC vs 传统 PD 自动舵部署对比  ·  船舶节能算法论文与部署案例综述",
          {"size": 20, "color": RGBColor(0x9F, 0xC5, 0xE8), "align": PP_ALIGN.CENTER})])
textbox(s, Inches(1), Inches(5.6), Inches(11.3), Inches(0.8),
        [("2026-08-07   |   数据均经两轮来源核对，出处见调研文档附录",
          {"size": 14, "color": GRAY, "align": PP_ALIGN.CENTER})])

# ============ S02 目录 ============
s = add_slide()
header(s, "汇报提纲")
bullets(s, Inches(1.0), Inches(1.6), Inches(11), Inches(5.5), [
    (0, "01  背景与法规驱动：EEXI/CII 机制详解，MPC 向船舶控制下沉", None, False),
    (0, "02  MPC vs PD 自动舵：原理、控制效果、证据强度、调试难度、实时性与稳定性", None, False),
    (0, "03  船舶节能五大算法路线：节能率、关键数据、优缺点", None, False),
    (0, "04  工业部署案例：规模化部署与实测细节", None, False),
    (0, "05  技术趋势判断与本项目落地建议", None, False),
], size=19)

# ============ S03 背景 ============
s = add_slide()
header(s, "背景：节能从「可选项」变为「合规刚需」", "法规 + 控制理论双重驱动")
bullets(s, Inches(0.6), Inches(1.7), Inches(6.4), Inches(5.4), [
    (0, "IMO EEXI / CII 强制生效（2023-01-01）", BLUE, True),
    (1, "EEXI：≥400 GT，一次性技术能效合规"),
    (1, "CII：≥5,000 GT，年度碳强度 A–E 评级"),
    (1, "连续 3 年 D 或 1 年 E → 必须提交纠正计划"),
    (1, "EU ETS 2024 年起纳入航运，碳成本显性化"),
    (0, "MPC 向船舶控制下沉", BLUE, True),
    (1, "嵌入式 QP 求解器成熟：水下机器人实测 7 ms @ 20 Hz（acados，JMSE 2024）"),
    (1, "船舶航向动态慢（控制周期 0.1–0.2 s），算力足够"),
    (1, "学术界大量仿真/船模成果，实船落地仍少"),
], size=15)
pic_fit(s, os.path.join(FIG, "fig6_architecture.png"), Inches(7.1), Inches(1.7),
        Inches(5.9), Inches(5.3))
textbox(s, Inches(7.1), Inches(6.9), Inches(5.9), Inches(0.4),
        [("船舶节能分层架构总览", {"size": 12, "color": GRAY, "align": PP_ALIGN.CENTER})])

# ============ S04 法规详解（新增） ============
s = add_slide()
header(s, "法规背景详解：EEXI 与 CII 的约束机制", "MARPOL 附则 VI 修正案（MEPC.328(76)），IMO 官方 FAQ 口径")
table(s, Inches(0.5), Inches(1.7), Inches(12.4), Inches(3.6), [
    ["项目", "EEXI（技术能效指数）", "CII（营运碳强度指标）"],
    ["适用船舶", "≥400 GT 各船型", "≥5,000 GT 各船型"],
    ["核心机制", "实际值 attained ≤ 要求值 required\n（required 按船型/吨位以折减系数确定）", "年度营运碳强度（AER，单位运输功 CO2）\n对照 G2 参考线 + G3 逐年折减系数"],
    ["合规方式", "一次性技术验证；\n常用轴/主机功率限制（ShaPoLi/EPL）", "按 G4 边界评 A–E 级；连续 3 年 D 或 1 年 E\n须在 SEEMP Part III 提交纠正措施计划"],
    ["关键时间线", "2022-11-01 生效\n2023-01-01 起须计算 attained EEXI", "2023 年起数据收集，2024 年首次评级\n法规要求 2026-01-01 前完成复审"],
], col_widths=[0.9, 2.4, 2.7], size=12)
bullets(s, Inches(0.6), Inches(5.6), Inches(12.2), Inches(1.7), [
    (0, "对节能技术选型的驱动：船东自 2022 年起集中评估船队并规划改装；EU ETS 自 2024 年纳入航运，碳成本显性化；", None, False),
    (0, "厂商已直接以「改善 CII 评级」作为产品卖点（Anschütz、Yara 等）——节能软件从「节油回本」变为「合规刚需」。", None, False),
], size=14)

# ============ S05 原理对比 ============
s = add_slide()
header(s, "MPC vs PD：原理差异", "前视预测 + 显式约束 + 多目标代价函数")
pic_fit(s, os.path.join(FIG, "fig2_mpc_principle.png"), Inches(0.4), Inches(1.7),
        Inches(7.3), Inches(5.2))
bullets(s, Inches(8.0), Inches(1.8), Inches(5.0), Inches(5.4), [
    (0, "PD/PID 自动舵", RED, True),
    (1, "δ = Kp·e + Kd·r：看着误差打舵"),
    (1, "无模型、几次乘加、单片机可跑"),
    (1, "约束靠事后饱和截断"),
    (0, "MPC 自动舵", BLUE, True),
    (1, "每周期解一个有限时域优化问题"),
    (1, "用 Nomoto 模型预测未来 Np 步响应"),
    (1, "舵角 |δ|≤35°、舵速 ≤3°/s 显式入约束"),
    (1, "代价函数可计入操舵能量（省油）"),
], size=15)

# ============ S06 效果对比 ============
s = add_slide()
header(s, "控制效果：仿真/船模层面 MPC 全面占优", "注意：尚无全尺度实船 MPC vs PID 公开定量对比")
pic_fit(s, os.path.join(FIG, "fig1_step_response.png"), Inches(0.4), Inches(1.7),
        Inches(7.3), Inches(4.6))
table(s, Inches(7.9), Inches(1.7), Inches(5.1), Inches(4.4), [
    ["研究", "关键结果"],
    ["Jannaty 2023\n(护卫舰仿真)", "到达目标航向：\nMPC 4 s vs PID 15 s"],
    ["Zhang 2025\n(内河船 NMPC)", "平均横偏误差降 82%\n(9.497→1.731 m)"],
    ["Appl. Sci. 2026\n(桥区水域)", "调节时间 120–160 s\nvs PID >600 s"],
    ["RA-L 2017\n(波浪场机器人)", "位置误差降 74%\n(MPC vs PD，仿真)"],
], col_widths=[1.15, 1.5], size=11)
textbox(s, Inches(0.5), Inches(6.55), Inches(12.4), Inches(0.7),
        [("［证据边界］对比几乎全部来自仿真/水池船模；预测模型与仿真模型同源会高估 MPC 优势（He 2023 警示）；",
          {"size": 13, "color": RED}),
         ("Zhang 2025 中舵机控制投入指标（AACE）在直道/弯道工况 PID 反优——「占优」限跟踪精度维度",
          {"size": 13, "color": RED})])

# ============ S07 证据强度分级（新增） ============
s = add_slide()
header(s, "证据强度分级：本文数据的可靠性坐标", "引用前先看数据站在哪一级")
table(s, Inches(0.5), Inches(1.7), Inches(12.4), Inches(4.3), [
    ["证据等级", "代表数据", "可靠性评估"],
    ["实船试航 / 第三方实测", "ClassNK-NAPA 试航 3.8–3.9%；Yara 实测 3–5%；\nMalaysia Grace 3.4%；MRAS 全尺度 1–3%", "最高，可直接引用"],
    ["船模实验", "He 2023 MPC 路径跟踪（8 Hz 自由自航船模）", "较高，注意尺度效应"],
    ["试验校验的仿真", "Pelić 2023 减速 72–76%（船模+柴油机模型，偏差 <2%）", "中等，模型经试验校验"],
    ["纯仿真", "MPC vs PID 多数对比（4 s vs 15 s、74%）；\nMPC-EMS 节油 4.85–17.2%", "参考级，需实船/船模验证"],
    ["厂商宣称", "ABB ≤9%；Anschütz ≤5%；Yara 5–15%", "注意营销偏差，以第三方复核为准"],
], col_widths=[1.5, 3.4, 1.6], size=12)
bullets(s, Inches(0.6), Inches(6.25), Inches(12.2), Inches(1.1), [
    (0, "关键提醒：MPC vs PID 无全尺度实船公开对比；「预测模型与仿真模型同源会高估优势」（He 2023）；Li 2022（2.1–5.2%）为二手转引。", RED, False),
], size=14)

# ============ S08 调试难度 ============
s = add_slide()
header(s, "调试难度：MPC 上船的最大障碍", "「模型参数确定是实现 MPC 的最大障碍」—— He et al., 2023 原文")
pic_fit(s, os.path.join(FIG, "fig3_radar.png"), Inches(0.4), Inches(1.6),
        Inches(6.0), Inches(5.5))
bullets(s, Inches(6.8), Inches(1.7), Inches(6.2), Inches(5.6), [
    (0, "PID 侧：规则成熟、船员可整定", RED, True),
    (1, "Z-N 临界比例度法（无模型，步骤机械）"),
    (1, "Nomoto 极点配置解析公式：Kp=ωn²T/K，Kd=(2ζωnT−1)/K"),
    (0, "MPC 侧：模型 + 多参数双重门槛", BLUE, True),
    (1, "需 Z 形/回转试验做系统辨识，模型随装载/污底漂移"),
    (1, "Np、Nc、Q、R、采样周期、约束边界多参数耦合"),
    (1, "权重整定多靠 trial-and-error（DTU 2025 原文佐证）"),
    (1, "稳定性需专门构造（终端约束 / Lyapunov 约束）"),
    (0, "注：雷达图为基于文献证据的主观打分示意", GRAY, False),
], size=14)

# ============ S09 实时性与稳定性（新增） ============
s = add_slide()
header(s, "MPC 工程化两个关键问题：实时性与稳定性", "均有文献实测数据支撑")
bullets(s, Inches(0.6), Inches(1.7), Inches(6.2), Inches(5.6), [
    (0, "实时性：船舶慢动态下可满足", BLUE, True),
    (1, "船舶航向控制周期 0.1–0.2 s 足够（He 2023 船模 8 Hz）"),
    (1, "acados 实测：预测时域 60、20 Hz 下平均求解 7 ms，仅占采样周期 14%（Hu 2024，JMSE）"),
    (1, "嵌入式平台（x86/ARM）QP 求解器基准：HPIPM 长时域最快；OSQP 热启动稳定可用（arXiv 2510.21773）"),
    (1, "可选求解器：OSQP / qpSWIFT / acados / HPIPM"),
    (1, "自主拖船 MPC 采样 0.2 s（You et al., 2024，UCL）"),
], size=14)
bullets(s, Inches(7.0), Inches(1.7), Inches(6.0), Inches(5.6), [
    (0, "稳定性：需专门构造，不能默认获得", ORANGE, True),
    (1, "PID：Nomoto 闭环分析成熟，频域裕度/极点配置直观"),
    (1, "MPC：有限时域优化不天然保证闭环稳定"),
    (1, "经典框架：终端约束集 + 终端代价（Mayne et al., 2000, Automatica）"),
    (1, "船舶落地形式：Lyapunov 约束嵌入优化问题"),
    (1, "Gong et al., 2021（AUV，Ocean Engineering）"),
    (1, "Zhang et al., 2023（DP 船，JMSE 11(2):281）"),
], size=14)

# ============ S10 优缺点小结 ============
s = add_slide()
header(s, "MPC vs PD 部署优缺点小结")
table(s, Inches(0.5), Inches(1.6), Inches(12.4), Inches(3.9), [
    ["维度", "传统 PD/PID", "MPC"],
    ["控制效果", "无前瞻、海浪下易无效操舵", "响应快、超调小、抗扰强、特定工况舵动作更少"],
    ["调试难度", "整定规则成熟，船员可现场整定", "模型辨识是硬门槛；权重多靠 trial-and-error"],
    ["实时计算", "几次乘加，单片机即可", "每周期解 QP/NLP；船舶慢动态下可行（实测 7 ms）"],
    ["约束/安全", "事后饱和截断，约束下退化", "舵角/舵速约束显式满足；稳定性需专门构造"],
    ["认证与维护", "工业生态与船级社认证完备", "认证路径不成熟；需控制工程师维护"],
], col_widths=[0.9, 2.2, 2.6], size=13)
textbox(s, Inches(0.5), Inches(5.8), Inches(12.4), Inches(1.3),
        [("结论性判断：PD/PID 打底保证可靠性，MPC 作为性能增强层；", {"size": 17, "bold": True, "color": GREEN}),
         ("模型辨识与整定自动化（本仓库 auto_tuning 流程）是 MPC 能否落地的前置条件。",
          {"size": 17, "bold": True, "color": GREEN})])

# ============ S11 节能总览 ============
s = add_slide()
header(s, "船舶节能五大路线：量级差异极大", "来源：IMO / 船级社 / 实船试航，均为已核对数据")
pic_fit(s, os.path.join(FIG, "fig4_energy_measures.png"), Inches(0.5), Inches(1.7),
        Inches(12.3), Inches(5.6))

# ============ S12 节能方案优缺点（一） ============
s = add_slide()
header(s, "节能方案优缺点（一）：决策层", "节能大头所在")
table(s, Inches(0.5), Inches(1.6), Inches(12.4), Inches(5.3), [
    ["方案", "代表数据（已核对）", "优点", "缺点"],
    ["减速航行", "23→12 kn 每海里油耗降 72–76%（Pelić 2023）\nIMO：主机油耗潜力 10–50%，全船口径 3–12%", "杠杆最大；零硬件成本；直接改善 CII", "牺牲船期；受租约约束；低速偏离主机最佳负荷点"],
    ["气象航线优化", "IMO 口径 ≥3%，集装箱船可达 10%\nLi 2022 实船 2.1–5.2%（二手转引）", "节能+避离恶劣海况兼得；商用服务成熟", "依赖气象预报与增阻模型精度；收益随海况波动"],
    ["纵倾优化", "ClassNK-NAPA 实船试航：纵倾贡献 1.2%\n专项研究最多 4%；Eniram 实测 1–5%", "几乎零成本（调压载水）；实船背书强", "最优纵倾随工况变化；受稳性/装卸约束；幅度偏小"],
], col_widths=[1.0, 2.6, 1.8, 1.9], size=12)

# ============ S13 节能关键数据补充（新增） ============
s = add_slide()
header(s, "节能方案关键数据补充", "对比基线与口径已逐条核对，避免误读")
table(s, Inches(0.5), Inches(1.6), Inches(12.4), Inches(5.3), [
    ["方案", "补充数据（已核对）", "口径说明"],
    ["减速航行", "IMO 案例：56,000 DWT 敞舱口货船 13% 减速 → 日油耗 −34%；\n10% 减速 → 功率 −27%，计入航时后航次总节油约 19%", "特定船案例，非普适值；\n主机口径 10–50% vs 全船 3–12%"],
    ["航速+纵倾联合", "Du 2019（实船数据+ANN 模型）：动态纵倾 C1 省 4.96%/5.83%；\nANN 航速优化 C2 省 7.63%/7.57%；综合 C3 平均省 8.25%", "0.57%/3.69% 是相对三次方定律的\n额外节油，非总节油"],
    ["纵倾优化", "ClassNK 两次印度洋横渡专项：最优纵倾最多再省 4%", "与航速/航程优化收益可叠加"],
    ["MPC 能量管理", "Zhang 2022：双层 MPC 比传统单层 MPC 再省 17.2%（30 kW 柴电混合船）\n周妍 2024：比规则控制省 4.85%（客滚船）", "两者对比基线不同，不可互比；\n均为仿真"],
    ["气象航线", "算法基准研究：Isochrone / Isopone / DP / 3D-DP / Dijkstra\n五算法横向对比（Wang et al., Chalmers，北大西洋案例）", "仿真基准，无单一最优算法"],
], col_widths=[1.2, 4.0, 2.0], size=11.5)

# ============ S14 节能方案优缺点（二） ============
s = add_slide()
header(s, "节能方案优缺点（二）：操作层与控制层")
table(s, Inches(0.5), Inches(1.6), Inches(12.4), Inches(5.3), [
    ["方案", "代表数据（已核对）", "优点", "缺点"],
    ["MPC 能量管理\n(混合动力 EMS)", "双层 MPC 比传统 MPC 再省 17.2%（Zhang 2022）\n客滚船 MPC-EMS 节油 4.85%（周妍 2024）", "滚动时域天然适合负载预测；多目标可编码", "全部仿真验证无实船数据；仅限混合动力船型"],
    ["自动舵节能\n(减少无效操舵)", "IMO：0.1–1.0%；ABS：最高 1%、改造成本 $2 万\nMRAS 实船试验 1–3%（van Amerongen 经典）", "实施成本最低（软件级）；兼降舵机磨损", "单项幅度最小；收益依赖海况"],
    ["双舵 Toe Angle\n(舵系水动力优化)", "Anschütz 实测最大 4.7%、平均约 2%（宣称 ≤5%）\n自适应模式平均减少舵动作 25%", "双舵船额外收益；已有商用产品", "仅适用双舵船；厂商数据待第三方复核"],
], col_widths=[1.3, 2.6, 1.8, 1.9], size=12)

# ============ S15 部署案例（一） ============
s = add_slide()
header(s, "工业部署案例（一）：第三方实测集中在 3–10%", "厂商宣称值普遍高于实测，引用以船级社/试航数据为准")
pic_fit(s, os.path.join(FIG, "fig5_deploy_cases.png"), Inches(0.4), Inches(1.6),
        Inches(7.6), Inches(5.4))
bullets(s, Inches(8.2), Inches(1.7), Inches(4.9), Inches(5.6), [
    (0, "规模化部署", BLUE, True),
    (1, "ABB OCTOPUS：1000+ 船，回收期仅几个月"),
    (1, "DeepSea：EPS 全船队 300 艘，油耗预报误差 0.8%"),
    (1, "Orca AI：Seaspan 267+ 船，单船年省约 $10 万"),
    (0, "控制层现状", GREEN, True),
    (1, "节能自动舵 = 自适应+海浪滤波（Anschütz ECO）"),
    (1, "无公开的 MPC 实船自动舵部署案例"),
    (1, "印证：门槛不在算力，在模型获取与整定"),
], size=14)

# ============ S16 部署案例（二）（新增） ============
s = add_slide()
header(s, "工业部署案例（二）：实测细节与口径", "注意：各案例节油路径不同，不可直接横向比较")
table(s, Inches(0.5), Inches(1.6), Inches(12.4), Inches(4.6), [
    ["案例", "实测细节（已核对）", "节油路径"],
    ["Wärtsilä/Eniram\nVLCC 案例", "320,000+ DWT、450 天数据：年燃油成本 −2.6%\n（≈48.2 万美元 / 730 吨燃油）", "纵倾优化 + 污底评估"],
    ["Yara FuelOpt\n（独立复核）", "NAPA 独立分析确认 10–18%：\nSten Bothnia 12 个月 17.9%、Ekfjord 24 个月 10.3%", "推进功率闭环优化"],
    ["NAPA × Norsepower\n× SHI-ME", "转筒帆+航程优化六条航线平均 CO2 −19%\n（纽约–阿姆斯特丹线 28%，NAPA 贡献 10–12 个百分点）", "风助推 + 航程优化"],
    ["Orca AI\n（Seaspan）", "近距遭遇事件 −54%；2024 年合计减排约 19.5 万吨 CO2；\n单船年省约 $10 万 / 减 500 t CO2", "减少避碰机动与无谓变速"],
    ["Kongsberg\nEcoAdvisor", "DP 工况主机/推进器停机建议工具（DOF 试点）；\n官方未公布定量节油数据（流传的 8–10% 无法证实，已弃用）", "动力定位能耗优化"],
], col_widths=[1.5, 3.6, 1.6], size=11.5)
textbox(s, Inches(0.5), Inches(6.45), Inches(12.4), Inches(0.8),
        [("Anschütz 补充：自适应 ECO 模式平均减少舵动作约 25%（多年运营数据分析），案例折算节油约 4%。",
          {"size": 13, "color": GRAY})])

# ============ S17 技术趋势判断（新增） ============
s = add_slide()
header(s, "技术趋势判断", "基于两轮来源核对的四条判断")
bullets(s, Inches(0.8), Inches(1.8), Inches(11.8), Inches(5.4), [
    (0, "① 节能大头在决策层，控制层价值在「少损耗」", BLUE, True),
    (1, "自动舵贡献 0.1–1.0%（IMO 口径），但 7×24 在线、几乎零边际成本，兼降舵机磨损"),
    (0, "② MPC 可行域正在打开，瓶颈不是算力", BLUE, True),
    (1, "嵌入式求解器毫秒级（7 ms @ 20 Hz）；真正门槛是模型辨识与整定自动化"),
    (0, "③ 数据驱动自学习模型成为商用标配", BLUE, True),
    (1, "ClassNK-NAPA 油耗预测精度 99.6%；DeepSea 周油耗预报误差 0.8%；纯机理模型路线已少见"),
    (0, "④ 法规持续加码，合规刚需取代节油回本逻辑", BLUE, True),
    (1, "CII 纠正计划机制 + 2026 年法规复审可能进一步收紧；EU ETS 碳成本显性化"),
], size=16)

# ============ S18 落地建议 ============
s = add_slide()
header(s, "对本项目（auto_rudder）的落地建议")
table(s, Inches(0.5), Inches(1.6), Inches(12.4), Inches(4.2), [
    ["阶段", "目标", "关键动作"],
    ["短期\n1–3 个月", "PD 基线交付", "Nomoto 极点配置整定（Kp=ωn²T/K，ζ≈1）；\n海浪滤波+死区抑制无效操舵（对标 IMO 0.1–1.0% 收益）"],
    ["中期\n3–6 个月", "MPC 增强层开发", "攻克模型在线辨识与权重整定工具化；\n嵌入式求解器选型（acados/OSQP）；与 PD 基线 A/B 对比"],
    ["长期\n6–12 个月", "分层节能架构", "MPC 代价函数计入操舵能量，对接航速/纵倾优化指令；\n输出舵机磨损/能耗统计；Lyapunov-based MPC 保稳定性"],
], col_widths=[0.9, 1.4, 4.0], size=13)
textbox(s, Inches(0.5), Inches(6.1), Inches(12.4), Inches(1.0),
        [("核心逻辑：节能大头在决策层（气象航线 3–10%、减速航行主机口径可达 50%），控制层贡献 0.1–1% 但几乎零成本；",
          {"size": 15, "color": DARK}),
         ("先把 PD 整定与辨识流程工具化，再让 MPC 平滑接管性能增强。",
          {"size": 15, "color": DARK})])

# ============ S19 主要数据来源 ============
s = add_slide()
header(s, "主要数据来源", "完整 45 条参考文献见调研文档；所有定量数据经两轮原文核对")
bullets(s, Inches(0.6), Inches(1.7), Inches(6.1), Inches(5.5), [
    (0, "控制对比研究", BLUE, True),
    (1, "Jannaty et al., 2023, Kadikma 14(3)"),
    (1, "Zhang et al., 2025, Ocean Engineering 334:121592"),
    (1, "He et al., 2023, Ocean Engineering（船模实验）"),
    (1, "Zheng et al., 2025, JMSE 13(5):851（自动舵综述）"),
    (0, "整定 / 实时 / 稳定性", BLUE, True),
    (1, "Christensen et al., 2025, arXiv:2509.11235"),
    (1, "Hu et al., 2024, JMSE 12(1):94（acados 实测）"),
    (1, "Mayne et al., 2000, Automatica 36(6)"),
], size=13)
bullets(s, Inches(6.9), Inches(1.7), Inches(6.1), Inches(5.5), [
    (0, "节能算法", GREEN, True),
    (1, "Pelić et al., 2023, JMSE 11(3):675"),
    (1, "Du et al., 2019, TR-B 121:88–114"),
    (1, "IMO GreenVoyage2050 / GloMEEP 官方评估"),
    (1, "ClassNK-NAPA 全尺度试航（2014）"),
    (0, "部署案例与法规", GREEN, True),
    (1, "ABB / Wärtsilä-Eniram / Yara / Kongsberg / NAPA / DeepSea / Anschütz / Orca AI 官方资料"),
    (1, "IMO EEXI/CII 官方 FAQ（MEPC.328(76) 等）"),
], size=13)

# ============ S20 结尾 ============
s = add_slide()
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
textbox(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
        [("谢 谢", {"size": 48, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER})])
textbox(s, Inches(1), Inches(4.2), Inches(11.3), Inches(1.5),
        [("完整数据与 45 条参考文献：docs/ship_energy_research/MPC与船舶节能算法调研文档.md",
          {"size": 15, "color": RGBColor(0x9F, 0xC5, 0xE8), "align": PP_ALIGN.CENTER}),
         ("所有定量数据经两轮来源核对，修正记录见文档附录 A",
          {"size": 15, "color": RGBColor(0x9F, 0xC5, 0xE8), "align": PP_ALIGN.CENTER})])

# ============ 页码 ============
for idx, slide in enumerate(prs.slides, 1):
    if idx in (1, len(prs.slides._sldIdLst)):  # 跳过封面与结尾页
        continue
    textbox(slide, Inches(12.55), Inches(7.08), Inches(0.7), Inches(0.35),
            [(f"{idx} / {len(prs.slides._sldIdLst)}", {"size": 11, "color": GRAY, "align": PP_ALIGN.RIGHT})])

out = os.path.join(BASE, "船舶控制与节能算法调研汇报.pptx")
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
