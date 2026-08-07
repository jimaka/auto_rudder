# -*- coding: utf-8 -*-
"""将汇报 PPTX 高保真渲染为 PDF（与本仓库 make_ppt.py 的布局数据逐形状一致）。

方案：用 python-pptx 解析 pptx 的每个形状（位置/尺寸/颜色/字号/对齐），
再用 matplotlib 逐页绘制到多页 PDF；文字换行用 PIL 字体度量精确计算。
"""
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import Rectangle
import matplotlib.image as mpimg
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import ImageFont

BASE = os.path.dirname(os.path.abspath(__file__))
PPTX = os.path.join(BASE, "船舶控制与节能算法调研汇报.pptx")
PDF = os.path.join(BASE, "船舶控制与节能算法调研汇报.pdf")

REG = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
BOLD = "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc"
for fp in (REG, BOLD):
    fm.fontManager.addfont(fp)
plt.rcParams["font.family"] = "Noto Sans CJK JP"

EMU = 914400.0
_font_cache = {}

def pil_font(size_pt, bold):
    key = (round(size_pt * 4), bold)
    if key not in _font_cache:
        _font_cache[key] = ImageFont.truetype(BOLD if bold else REG, key[0])
    return _font_cache[key]

def text_w(text, size_pt, bold):
    return pil_font(size_pt, bold).getlength(text) / 4 / 72.0  # 英寸

def wrap_line(seg, size_pt, bold, max_w):
    """贪心换行：优先在空格处断行，CJK 逐字断行"""
    out, cur = [], ""
    for ch in seg:
        if text_w(cur + ch, size_pt, bold) <= max_w or not cur:
            cur += ch
        else:
            if ch == " ":
                out.append(cur); cur = ""
                continue
            sp = cur.rfind(" ")
            if sp > 0 and text_w(cur[:sp], size_pt, bold) > max_w * 0.4:
                out.append(cur[:sp]); cur = cur[sp + 1:] + ch
            else:
                out.append(cur); cur = ch
    out.append(cur)
    return out

def para_style(para):
    size, bold, color = 16.0, False, (0.1, 0.16, 0.29)
    for r in para.runs:
        if r.font.size: size = r.font.size.pt
        if r.font.bold is not None: bold = r.font.bold
        try:
            if r.font.color and r.font.color.rgb is not None:
                c = r.font.color.rgb
                color = (c[0] / 255, c[1] / 255, c[2] / 255)
        except Exception:
            pass
        break  # 本生成器每段单 run
    return size, bold, color

def render_textframe(ax, tf, x, y, w, h):
    """x,y 为左上角（英寸，pptx 坐标系）"""
    anchor = tf.vertical_anchor
    blocks = []
    for para in tf.paragraphs:
        text = "".join(r.text for r in para.runs)
        if not text.strip():
            continue
        size, bold, color = para_style(para)
        ls = para.line_spacing if para.line_spacing else 1.15
        sa = (para.space_after.pt if para.space_after else 6) / 72.0
        lines = []
        for seg in text.split("\n"):
            lines.extend(wrap_line(seg, size, bold, w - 0.1))
        lh = size * 1.22 * ls / 72.0
        blocks.append((lines, size, bold, color, lh, sa,
                       para.alignment or PP_ALIGN.LEFT))
    total_h = sum(len(b[0]) * b[4] + b[5] for b in blocks)
    if anchor == MSO_ANCHOR.MIDDLE:
        cy = y + max(0.0, (h - total_h) / 2)
    else:
        cy = y + 0.05
    for lines, size, bold, color, lh, sa, align in blocks:
        for ln in lines:
            lw = text_w(ln, size, bold)
            if align == PP_ALIGN.CENTER:
                lx = x + (w - lw) / 2
            elif align == PP_ALIGN.RIGHT:
                lx = x + w - lw - 0.05
            else:
                lx = x + 0.05
            ax.text(lx, cy + lh * 0.78, ln, fontsize=size,
                    fontweight="bold" if bold else "normal",
                    color=color, ha="left", va="baseline")
            cy += lh
        cy += sa

def cell_fill_rgb(cell):
    try:
        if cell.fill.type is not None and str(cell.fill.type) == "MSO_FILL_TYPE.BACKGROUND (5)":
            return None
        c = cell.fill.fore_color
        if c and c.rgb is not None:
            return (c.rgb[0] / 255, c.rgb[1] / 255, c.rgb[2] / 255)
    except Exception:
        pass
    return (1, 1, 1)

def render_table(ax, tbl_shape, x, y):
    t = tbl_shape.table
    col_ws = [c.width / EMU for c in t.columns]
    row_hs = [r.height / EMU for r in t.rows]
    cy = y
    for r, row in enumerate(t.rows):
        rh = row_hs[r]
        cx = x
        for c, cell in enumerate(row.cells):
            cw = col_ws[c]
            fill = cell_fill_rgb(cell)
            if fill is not None:
                ax.add_patch(Rectangle((cx, cy), cw, rh, facecolor=fill,
                                       edgecolor=(1, 1, 1), lw=0.8, zorder=2))
            render_textframe(ax, cell.text_frame, cx + 0.03, cy, cw - 0.06, rh)
            cx += cw
        cy += rh

def render_slide(fig, slide, sw_in, sh_in):
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, sw_in); ax.set_ylim(0, sh_in)
    ax.invert_yaxis()  # pptx 左上角原点
    ax.axis("off")
    ax.add_patch(Rectangle((0, 0), sw_in, sh_in, facecolor="white", zorder=0))
    for sh in slide.shapes:
        x, y = sh.left / EMU, sh.top / EMU
        w, h = sh.width / EMU, sh.height / EMU
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            import io
            img = mpimg.imread(io.BytesIO(sh.image.blob), format="png")
            ax.imshow(img, extent=[x, x + w, y + h, y], zorder=3,
                      interpolation="lanczos")
        elif sh.has_table:
            render_table(ax, sh, x, y)
        else:
            if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                try:
                    c = sh.fill.fore_color.rgb
                    ax.add_patch(Rectangle((x, y), w, h,
                                           facecolor=(c[0]/255, c[1]/255, c[2]/255),
                                           edgecolor="none", zorder=1))
                except Exception:
                    pass
            if sh.has_text_frame and sh.text_frame.text.strip():
                render_textframe(ax, sh.text_frame, x, y, w, h)

prs = Presentation(PPTX)
sw_in, sh_in = prs.slide_width / EMU, prs.slide_height / EMU
with PdfPages(PDF) as pdf:
    for slide in prs.slides:
        fig = plt.figure(figsize=(sw_in, sh_in))
        render_slide(fig, slide, sw_in, sh_in)
        pdf.savefig(fig, dpi=200)
        plt.close(fig)
print("saved:", PDF, "| pages:", len(prs.slides._sldIdLst))
